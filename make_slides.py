"""
make_slides.py
Generates a clean, modern PowerPoint presentation for the tissue inspection project.
Run from project root:
    py -3.11 make_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
ACCENT      = RGBColor(0x4F, 0x9C, 0xFF)   # bright blue
ACCENT2     = RGBColor(0x00, 0xD4, 0xAA)   # teal green
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG     = RGBColor(0x1A, 0x1A, 0x2E)   # slightly lighter dark
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW      = RGBColor(0xF3, 0x9C, 0x12)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helper functions ───────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def bg(slide, color=DARK_BG):
    add_rect(slide, 0, 0, W, H, color)

def accent_bar(slide, y=Inches(0.55), h=Inches(0.06)):
    add_rect(slide, 0, y, W, h, ACCENT)

def slide_title(slide, title, subtitle=None):
    add_text(slide, title, Inches(0.6), Inches(0.18), Inches(12), Inches(0.55),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(0.72), Inches(12), Inches(0.35),
                 size=14, color=LIGHT_GRAY)

def add_image_safe(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        box = add_rect(slide, x, y, w, h, CARD_BG)
        add_text(slide, f"[image: {os.path.basename(path)}]",
                 x + Inches(0.1), y + h/2 - Pt(10),
                 w - Inches(0.2), Pt(30),
                 size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, Inches(0.5), ACCENT)
add_rect(s, 0, H - Inches(0.5), W, Inches(0.5), ACCENT)

add_text(s, "ML-Based Visual Quality", Inches(1), Inches(1.3), Inches(11.3), Inches(1.1),
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Inspection System", Inches(1), Inches(2.3), Inches(11.3), Inches(1.0),
         size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(s, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06), ACCENT2)

add_text(s, "PDE4444 — Machine Learning for Engineers", Inches(1), Inches(3.75),
         Inches(11.3), Inches(0.45), size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "William Kojumian  ·  M01094013  ·  Middlesex University Dubai",
         Inches(1), Inches(4.25), Inches(11.3), Inches(0.4),
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Engineering Problem", "What are we solving?")

# Pipeline boxes
labels = ["Camera", "OpenCV", "Preprocess\n(Greyscale\n224×224)", "CNN Model", "PASS / FAIL"]
colors = [CARD_BG, CARD_BG, CARD_BG, ACCENT, ACCENT2]
bw, bh = Inches(1.9), Inches(1.4)
start_x = Inches(0.5)
y = Inches(2.4)
gap = Inches(0.22)

for i, (lbl, col) in enumerate(zip(labels, colors)):
    x = start_x + i * (bw + gap)
    add_rect(s, x, y, bw, bh, col)
    add_text(s, lbl, x, y + Inches(0.25), bw, bh - Inches(0.25),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(labels) - 1:
        add_text(s, "→", x + bw + Inches(0.02), y + Inches(0.45),
                 gap + Inches(0.05), Inches(0.5), size=18, bold=True,
                 color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "A tissue box is placed in the inspection area. The camera captures an image,\n"
            "the model classifies it as PASS (undamaged) or FAIL (defective) in real time.",
         Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8),
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Defect types
add_text(s, "Defect Types:", Inches(0.6), Inches(5.15), Inches(3), Inches(0.4),
         size=13, bold=True, color=ACCENT)
defects = ["Punch-face denting", "Corner crushing", "Surface scratching"]
for i, d in enumerate(defects):
    add_text(s, f"  {i+1}. {d}", Inches(0.6), Inches(5.55 + i*0.35), Inches(4), Inches(0.35),
             size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Dataset", "824 images collected under controlled conditions")

# Stats cards
stats = [("824", "Total Images"), ("313", "PASS"), ("511", "FAIL"), ("70/15/15", "Train/Val/Test %")]
scol  = [ACCENT, GREEN, RED, ACCENT2]
cw, ch = Inches(2.8), Inches(1.5)
sx = Inches(0.5)
for i, ((num, lbl), col) in enumerate(zip(stats, scol)):
    x = sx + i * (cw + Inches(0.25))
    add_rect(s, x, Inches(1.1), cw, ch, CARD_BG)
    add_rect(s, x, Inches(1.1), cw, Inches(0.06), col)
    add_text(s, num, x, Inches(1.2), cw, Inches(0.75),
             size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x, Inches(1.9), cw, Inches(0.4),
             size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Sample images
img_base = "report_v2"
imgs = [
    ("original__IMG_9928.jpg",  "PASS — Original"),
    ("grayscale__IMG_9928.jpg", "PASS — Greyscale"),
    ("original__IMG_0020.jpg",  "FAIL — Original"),
    ("grayscale__IMG_0020.jpg", "FAIL — Greyscale"),
]
iw, ih = Inches(2.7), Inches(2.7)
ix_start = Inches(0.5)
for i, (fname, caption) in enumerate(imgs):
    ix = ix_start + i * (iw + Inches(0.3))
    add_image_safe(s, os.path.join(img_base, fname), ix, Inches(2.9), iw, ih)
    add_text(s, caption, ix, Inches(5.65), iw, Inches(0.4),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why Greyscale
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Why Greyscale?", "A deliberate design decision")

# Two columns
add_rect(s, Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5), CARD_BG)
add_rect(s, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.06), RED)
add_text(s, "Without Greyscale (RGB)", Inches(0.6), Inches(1.15), Inches(5.6), Inches(0.5),
         size=15, bold=True, color=RED)
rgb_points = [
    "Model sees colour differences between boxes",
    "May learn: white box = PASS, beige box = FAIL",
    "That's colour detection, not damage detection",
    "Will fail on new box colours it hasn't seen",
]
for i, pt in enumerate(rgb_points):
    add_text(s, f"✗  {pt}", Inches(0.7), Inches(1.85 + i*0.7), Inches(5.4), Inches(0.6),
             size=13, color=LIGHT_GRAY)

add_rect(s, Inches(6.9), Inches(1.1), Inches(5.8), Inches(5.5), CARD_BG)
add_rect(s, Inches(6.9), Inches(1.1), Inches(5.8), Inches(0.06), GREEN)
add_text(s, "With Greyscale", Inches(7.0), Inches(1.15), Inches(5.6), Inches(0.5),
         size=15, bold=True, color=GREEN)
gs_points = [
    "Colour information completely removed",
    "Model forced to learn shape, texture, deformation",
    "Works on any box colour, any brand",
    "94.4% accuracy — genuine structural damage detection",
]
for i, pt in enumerate(gs_points):
    add_text(s, f"✓  {pt}", Inches(7.0), Inches(1.85 + i*0.7), Inches(5.6), Inches(0.6),
             size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CNN Architectures
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "CNN Architectures", "Custom-built vs Transfer Learning")

# Custom CNN
add_rect(s, Inches(0.4), Inches(1.1), Inches(5.9), Inches(5.8), CARD_BG)
add_rect(s, Inches(0.4), Inches(1.1), Inches(5.9), Inches(0.06), ACCENT)
add_text(s, "Custom CNN (from scratch)", Inches(0.5), Inches(1.15),
         Inches(5.7), Inches(0.45), size=14, bold=True, color=ACCENT)

layers = [
    ("Input",    "224×224×3 greyscale"),
    ("Block 1",  "Conv(3→32) → BN → ReLU → MaxPool  →  112×112×32"),
    ("Block 2",  "Conv(32→64) → BN → ReLU → MaxPool  →  56×56×64"),
    ("Block 3",  "Conv(64→128) → BN → ReLU → MaxPool  →  28×28×128"),
    ("Block 4",  "Conv(128→256) → BN → ReLU → AvgPool  →  1×1×256"),
    ("Head",     "Linear(256→128) → ReLU → Dropout(0.5) → Linear(128→2)"),
    ("Output",   "Softmax → PASS / FAIL"),
]
lcolors = [ACCENT2, ACCENT, ACCENT, ACCENT, ACCENT, YELLOW, GREEN]
for i, ((lname, ldesc), lc) in enumerate(zip(layers, lcolors)):
    y = Inches(1.8 + i * 0.65)
    add_rect(s, Inches(0.55), y, Inches(1.0), Inches(0.45), lc)
    add_text(s, lname, Inches(0.55), y, Inches(1.0), Inches(0.45),
             size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s, ldesc, Inches(1.65), y + Inches(0.06), Inches(4.5), Inches(0.38),
             size=10, color=LIGHT_GRAY)

add_text(s, "~0.8M parameters", Inches(0.5), Inches(6.5), Inches(5.7), Inches(0.35),
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# MobileNetV2
add_rect(s, Inches(6.8), Inches(1.1), Inches(5.9), Inches(5.8), CARD_BG)
add_rect(s, Inches(6.8), Inches(1.1), Inches(5.9), Inches(0.06), ACCENT2)
add_text(s, "MobileNetV2 (Transfer Learning)", Inches(6.9), Inches(1.15),
         Inches(5.7), Inches(0.45), size=14, bold=True, color=ACCENT2)

mob_points = [
    ("Pretrained on ImageNet (1.2M images)", LIGHT_GRAY),
    ("18 inverted residual bottleneck blocks", LIGHT_GRAY),
    ("13 blocks FROZEN  —  use pretrained knowledge", ACCENT),
    ("Last 5 blocks FINE-TUNED  —  adapt to our data", ACCENT2),
    ("Global Average Pooling  →  1,280-dim vector", LIGHT_GRAY),
    ("Dropout(0.3) → Linear(1280→2) → Softmax", LIGHT_GRAY),
]
for i, (pt, col) in enumerate(mob_points):
    add_text(s, f"→  {pt}", Inches(6.9), Inches(1.85 + i*0.7),
             Inches(5.6), Inches(0.6), size=12, color=col)

add_text(s, "~3.4M parameters  ·  Best Val Accuracy: 96.8%",
         Inches(6.8), Inches(6.5), Inches(5.9), Inches(0.35),
         size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Activation & Optimiser Comparison
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Activation Functions & Optimisers", "Finding the best combination")

# Activation table
add_text(s, "Activation + Optimiser Comparison", Inches(0.5), Inches(1.1),
         Inches(6), Inches(0.4), size=14, bold=True, color=ACCENT)

headers = ["Activation", "Optimiser", "Val Acc", ""]
col_x   = [Inches(0.5), Inches(2.3), Inches(4.1), Inches(5.3)]
col_w   = [Inches(1.7), Inches(1.7), Inches(1.1), Inches(1.2)]

for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s, x, Inches(1.55), w, Inches(0.38), ACCENT)
    add_text(s, h, x, Inches(1.55), w, Inches(0.38),
             size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

rows = [
    ("ReLU", "Adam", "95.2%", "★ BEST"),
    ("ReLU", "SGD",  "90.3%", ""),
    ("GELU", "Adam", "94.4%", ""),
    ("GELU", "SGD",  "88.7%", ""),
]
row_colors = [GREEN, CARD_BG, CARD_BG, CARD_BG]
for j, (row, rc) in enumerate(zip(rows, row_colors)):
    ry = Inches(1.93 + j * 0.45)
    add_rect(s, Inches(0.5), ry, Inches(6.0), Inches(0.42), rc)
    for i, (val, x, w) in enumerate(zip(row, col_x, col_w)):
        tc = ACCENT if (j == 0 and i == 3) else WHITE if j == 0 else LIGHT_GRAY
        add_text(s, val, x, ry, w, Inches(0.42),
                 size=12, color=tc, align=PP_ALIGN.CENTER, bold=(j==0))

# Optimiser table
add_text(s, "Optimisation Methods", Inches(0.5), Inches(4.1),
         Inches(6), Inches(0.4), size=14, bold=True, color=ACCENT2)

opt_rows = [
    ("Zero-order",  "Random hyperparameter search",    "No gradients"),
    ("First-order", "SGD with Momentum",               "∇L, fixed LR"),
    ("First-order", "Adam (adaptive)",                 "∇L, adaptive LR — BEST"),
    ("Second-order","L-BFGS (Hessian approx)",         "Slow, impractical for mini-batch"),
]
for j, (order, name, note) in enumerate(opt_rows):
    ry = Inches(4.6 + j * 0.52)
    add_rect(s, Inches(0.5), ry, Inches(1.5), Inches(0.45), ACCENT if j == 2 else CARD_BG)
    add_text(s, order, Inches(0.5), ry, Inches(1.5), Inches(0.45),
             size=10, color=WHITE if j == 2 else LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, name,  Inches(2.1), ry + Inches(0.04), Inches(2.5), Inches(0.38),
             size=11, bold=(j==2), color=WHITE if j == 2 else LIGHT_GRAY)
    add_text(s, note,  Inches(4.7), ry + Inches(0.04), Inches(2.3), Inches(0.38),
             size=10, color=ACCENT2 if j==2 else LIGHT_GRAY)

# Comparison image
add_image_safe(s, "report_v2/activation_optimizer_comparison.png",
               Inches(6.8), Inches(1.1), Inches(5.9), Inches(5.7))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Results
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Model Results", "Greyscale training — colour-invariant damage detection")

# Big accuracy cards
models_data = [
    ("Logistic\nRegression\n(PCA-50)", "68.5%", "66.8%", RED,   "Linear baseline\nCannot learn\nnon-linear damage"),
    ("Custom CNN\n(from scratch)",      "87.9%", "86.4%", YELLOW,"4-block architecture\nLearns features\nautomatically"),
    ("MobileNetV2\n(Transfer)",         "94.4%", "92.8%", GREEN, "Best model\nPretrained ImageNet\nFine-tuned"),
]
cw, ch = Inches(3.8), Inches(4.8)
for i, (name, acc, f1, col, note) in enumerate(models_data):
    x = Inches(0.4) + i * (cw + Inches(0.3))
    add_rect(s, x, Inches(1.05), cw, ch, CARD_BG)
    add_rect(s, x, Inches(1.05), cw, Inches(0.07), col)
    add_text(s, name, x, Inches(1.15), cw, Inches(0.9),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, acc, x, Inches(2.1), cw, Inches(1.0),
             size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, f"F1: {f1}", x, Inches(3.1), cw, Inches(0.4),
             size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(0.3), Inches(3.6), cw - Inches(0.6), Inches(0.04), col)
    add_text(s, note, x, Inches(3.75), cw, Inches(1.1),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "68.5%  →  87.9%  →  94.4%  —  Each model outperforms the last",
         Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.45),
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Confusion Matrix + PCA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Evaluation — Test Set (124 images)", "MobileNetV2 greyscale model")

add_image_safe(s, "report_v2/confusion_matrix.png",
               Inches(0.4), Inches(1.1), Inches(5.5), Inches(4.5))
add_image_safe(s, "report_v2/pca_features.png",
               Inches(6.3), Inches(1.1), Inches(6.0), Inches(4.5))

# Metrics row
metrics = [("94.4%", "Accuracy"), ("90.0%", "Precision"), ("95.7%", "Recall"), ("92.8%", "F1 Score"), ("7/124", "Errors")]
mw = Inches(2.3)
for i, (val, lbl) in enumerate(metrics):
    x = Inches(0.4) + i * (mw + Inches(0.12))
    add_rect(s, x, Inches(5.8), mw, Inches(0.9), CARD_BG)
    add_text(s, val, x, Inches(5.8), mw, Inches(0.5),
             size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x, Inches(6.25), mw, Inches(0.35),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "Clean class separation in PCA confirms model learned structural damage, not colour patterns",
         Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.4),
         size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Training Curves
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Training Convergence", "MobileNetV2 — 30 epochs, ReLU + Adam")

add_image_safe(s, "report_v2/training_curves.png",
               Inches(1.0), Inches(1.1), Inches(11.3), Inches(4.8))

points = [
    ("Val accuracy reaches 96.8% — model saved at peak"),
    ("Train and val curves track closely — no overfitting"),
    ("Stabilises after epoch 20 — convergence confirmed"),
]
for i, pt in enumerate(points):
    add_rect(s, Inches(0.5), Inches(6.1 + i*0.38), Inches(0.3), Inches(0.28), ACCENT2)
    add_text(s, pt, Inches(0.9), Inches(6.08 + i*0.38), Inches(11.5), Inches(0.35),
             size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Live Demo Results
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_title(s, "Live Demo — Real-World Videos", "New unseen data — iPhone camera")

# FAIL card
add_rect(s, Inches(0.4), Inches(1.1), Inches(5.9), Inches(5.5), CARD_BG)
add_rect(s, Inches(0.4), Inches(1.1), Inches(5.9), Inches(0.06), RED)
add_text(s, "IMG_0931 — DAMAGED BOX", Inches(0.5), Inches(1.2),
         Inches(5.7), Inches(0.45), size=14, bold=True, color=RED)
fail_stats = [
    ("VERDICT", "FAIL — DAMAGED", RED),
    ("FAIL frames", "87.0%  (40/46)", LIGHT_GRAY),
    ("Avg Confidence", "87.2%", LIGHT_GRAY),
    ("Defect", "Dented corners", LIGHT_GRAY),
]
for i, (lbl, val, col) in enumerate(fail_stats):
    add_text(s, lbl, Inches(0.6), Inches(1.85 + i*0.75), Inches(2.2), Inches(0.55),
             size=11, color=LIGHT_GRAY)
    add_text(s, val, Inches(2.9), Inches(1.85 + i*0.75), Inches(3.2), Inches(0.55),
             size=14, bold=True, color=col)
add_rect(s, Inches(0.8), Inches(4.8), Inches(4.8), Inches(0.06), RED)
add_text(s, "Model correctly identified structural damage\nwith high confidence across almost all frames",
         Inches(0.5), Inches(5.0), Inches(5.7), Inches(0.7),
         size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

# PASS card
add_rect(s, Inches(7.0), Inches(1.1), Inches(5.9), Inches(5.5), CARD_BG)
add_rect(s, Inches(7.0), Inches(1.1), Inches(5.9), Inches(0.06), GREEN)
add_text(s, "IMG_0932 — UNDAMAGED BOX", Inches(7.1), Inches(1.2),
         Inches(5.7), Inches(0.45), size=14, bold=True, color=GREEN)
pass_stats = [
    ("VERDICT", "PASS — GOOD CONDITION", GREEN),
    ("PASS frames", "84.0%  (42/50)", LIGHT_GRAY),
    ("Avg Confidence", "87.8%", LIGHT_GRAY),
    ("Defect", "None", LIGHT_GRAY),
]
for i, (lbl, val, col) in enumerate(pass_stats):
    add_text(s, lbl, Inches(7.2), Inches(1.85 + i*0.75), Inches(2.2), Inches(0.55),
             size=11, color=LIGHT_GRAY)
    add_text(s, val, Inches(9.5), Inches(1.85 + i*0.75), Inches(3.2), Inches(0.55),
             size=14, bold=True, color=col)
add_rect(s, Inches(7.3), Inches(4.8), Inches(4.8), Inches(0.06), GREEN)
add_text(s, "Model correctly cleared the undamaged box.\nMinority of frames uncertain due to motion blur.",
         Inches(7.1), Inches(5.0), Inches(5.7), Inches(0.7),
         size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "Both videos were brand new — never seen during training",
         Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.4),
         size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, H, DARK_BG)
accent_bar(s)
slide_title(s, "Conclusion", "")

points = [
    (ACCENT,  "94.4% test accuracy",           "on greyscale images — colour-invariant damage detection"),
    (GREEN,   "3-model comparison",             "Logistic Regression → Custom CNN → MobileNetV2"),
    (ACCENT2, "Greyscale preprocessing",        "deliberate design choice — forces structural feature learning"),
    (YELLOW,  "Real-world video validation",    "both test videos correctly classified with 87%+ confidence"),
    (ACCENT,  "Full experimental rigor",        "train/val/test split, 5-fold cross-validation, overfitting analysis"),
]
for i, (col, title, desc) in enumerate(points):
    y = Inches(1.4 + i * 0.95)
    add_rect(s, Inches(0.5), y, Inches(0.08), Inches(0.6), col)
    add_text(s, title, Inches(0.75), y, Inches(3.5), Inches(0.6),
             size=15, bold=True, color=col)
    add_text(s, desc,  Inches(4.4), y + Inches(0.08), Inches(8.3), Inches(0.5),
             size=13, color=LIGHT_GRAY)

add_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.06), ACCENT)
add_text(s, "William Kojumian  ·  M01094013  ·  PDE4444  ·  Middlesex University Dubai",
         Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
         size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "report_v2/PDE4444_Tissue_Inspection_Presentation.pptx"
os.makedirs("report_v2", exist_ok=True)
prs.save(out)
print(f"[done] Saved: {out}")
